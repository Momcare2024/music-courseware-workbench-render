from __future__ import annotations

from pathlib import Path


def export_image_deck(project: dict, base_dir: Path, output_path: Path) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx 依赖，请先安装 requirements.txt") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    images_dir = base_dir / "artifacts" / "images"

    pages = project.get("pages") or []
    if not pages:
        raise RuntimeError("当前项目还没有页面")

    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        image_file = str(page.get("image_file") or "")
        image_path = images_dir / image_file
        if image_file and image_path.exists():
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        else:
            textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.2))
            frame = textbox.text_frame
            frame.text = f"{page.get('title') or '未生成页面'}：尚未生成图片"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
